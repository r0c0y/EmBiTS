import os, hashlib, io, tempfile, re
from datetime import datetime
from PIL import Image
import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from ocr_engine import get_ocr_manager, OCRResult
from ocr_storage import store_ocr_result

EXTS = {".txt",".pdf",".docx",".xlsx",".xls",".pptx",".png",".jpg",".jpeg",".tiff",".bmp",".csv",".md",".json"}

_DATE_PATTERNS = [
    (r"\b(\d{4}-\d{2}-\d{2})\b", "%Y-%m-%d"),
    (r"\b(\d{1,2}[/\-\.]\d{1,2}[/\-\.]\d{4})\b", None),
    (r"\b(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]* \d{1,2},? \d{4}\b", "%b %d, %Y"),
    (r"\b\d{1,2} (Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]* \d{4}\b", "%d %b %Y"),
]

def detect_date(text):
    text = text or ""
    for pat, fmt in _DATE_PATTERNS:
        m = re.search(pat, text, re.IGNORECASE)
        if m:
            d = m.group(0) if fmt is None else m.group(1)
            for try_fmt in [fmt] if fmt else ["%d-%m-%Y","%m-%d-%Y","%d/%m/%Y","%m/%d/%Y"]:
                try:
                    return datetime.strptime(d, try_fmt).strftime("%Y-%m-%d")
                except: pass
    return None

def chunk_text(text, size=800, overlap=200):
    text = text or ""
    if not text.strip():
        return [""]

    # Detect form/template structure: repeated "Field Name:" or "Field Name :" patterns
    form_lines = text.split("\n")
    form_field_count = sum(1 for l in form_lines if re.match(r'^[A-Za-z][A-Za-z /-]+[?:]\s*$', l.strip()) or re.match(r'^[A-Za-z][A-Za-z /-]+[?:].{0,20}$', l.strip()))
    is_form_like = form_field_count >= 3 and len(text) < size * 3

    # For form-like documents, keep everything as a single chunk (short templates)
    if is_form_like:
        return [text.strip()]

    # Split on paragraph boundaries (double newlines) first
    paragraphs = re.split(r'\n\s*\n', text)
    paragraphs = [p.strip() for p in paragraphs if p.strip()]

    if not paragraphs:
        return [""]

    chunks = []
    current = ""
    for p in paragraphs:
        candidate = (current + "\n\n" + p).strip() if current else p
        if len(candidate) <= size:
            current = candidate
        else:
            if current:
                chunks.append(current)
            if len(p) > size:
                sentences = re.split(r'(?<=[.!?])\s+', p)
                if len(sentences) > 1:
                    cur_s = ""
                    for s in sentences:
                        cand = (cur_s + " " + s).strip() if cur_s else s
                        if len(cand) <= size:
                            cur_s = cand
                        else:
                            if cur_s:
                                chunks.append(cur_s)
                            cur_s = s
                    if cur_s:
                        chunks.append(cur_s)
                else:
                    for i in range(0, len(p), size - overlap):
                        chunks.append(p[i:i+size].strip())
                        if i + size >= len(p):
                            break
            else:
                current = p
    if current:
        chunks.append(current)

    return chunks or [""]

def content_hash(raw):
    return hashlib.sha256(raw).hexdigest()

def _extract_structured(f, ext):
    f.seek(0)
    if ext in (".txt",".csv",".md",".json"):
        return f.read().decode("utf-8", errors="ignore")
    if ext == ".docx":
        import zipfile
        doc = Document(f)
        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        for table in doc.tables:
            if not table.rows:
                continue
            table_lines = []
            col_count = len(table.rows[0].cells)
            for r_idx, row in enumerate(table.rows):
                cells_text = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                row_line = "| " + " | ".join(cells_text) + " |"
                table_lines.append(row_line)
                if r_idx == 0:
                    sep = "|" + "|".join("---" for _ in range(col_count)) + "|"
                    table_lines.append(sep)
            parts.append("\n" + "\n".join(table_lines) + "\n")
        try:
            f.seek(0)
            with zipfile.ZipFile(f) as z:
                media_files = [name for name in z.namelist() if name.startswith("word/media/")]
                if media_files:
                    mgr = get_ocr_manager()
                    for m_file in media_files:
                        img_data = z.read(m_file)
                        img = Image.open(io.BytesIO(img_data))
                        res = mgr.process_image(img)
                        if res.text.strip():
                            parts.append(f"\n[Embedded Image: {m_file.split('/')[-1]}]\n{res.text}")
        except Exception:
            pass
        return "\n".join(parts)
    if ext in (".xlsx",".xls"):
        wb = load_workbook(f, read_only=True, data_only=True)
        t = []
        for s in wb.worksheets:
            all_rows = [list(r) for r in s.iter_rows(values_only=True)]
            if not all_rows:
                continue
            col_count = max(len(r) for r in all_rows) if all_rows else 0
            if col_count == 0:
                continue
            header_sep = False
            for r_idx, row_data in enumerate(all_rows):
                row_text = [str(c) if c is not None else "" for c in row_data]
                while len(row_text) < col_count:
                    row_text.append("")
                cleaned = [c.replace("\n", " ").strip() for c in row_text]
                line = "| " + " | ".join(cleaned) + " |"
                t.append(line)
                if r_idx == 0:
                    sep = "|" + "|".join("---" for _ in range(col_count)) + "|"
                    t.append(sep)
        return "\n".join(t)
    if ext == ".pptx":
        prs = Presentation(f)
        return "\n".join(sh.text for slide in prs.slides for sh in slide.shapes if hasattr(sh,"text") and sh.text.strip())
    return None

def _save_temp(file_obj, suffix):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    try:
        file_obj.seek(0)
        tmp.write(file_obj.read())
        tmp.close()
        return tmp.name
    except Exception:
        try: tmp.close()
        except: pass
        return None

def ingest(file_path, filename, size=0, base_dir=None):
    base_dir = base_dir or os.path.dirname(os.path.abspath(__file__))
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in EXTS:
        raise ValueError(f"Unsupported format: {ext}")

    structured = None
    # For text-based formats (txt, docx, xlsx, pptx, csv, md, json), read as file object
    if ext in (".txt", ".csv", ".md", ".json", ".docx", ".xlsx", ".xls", ".pptx"):
        with open(file_path, "rb") as f:
            structured = _extract_structured(f, ext)
    
    if structured is not None:
        doc_date = detect_date(structured)
        return {
            "text": structured,
            "detected_date": doc_date,
            "source_type": ext.lstrip("."),
            "meta": {"engine": "native"},
            "ocr_json": None,
            "ocr_markdown": None,
            "ocr_engine": "native",
            "page_count": 1,
        }

    mgr = get_ocr_manager()
    try:
        ocr_result = mgr.process_with_fallback(file_path, ext)
    except Exception as e:
        raise ValueError(f"OCR processing failed: {e}") from e

    if not ocr_result.text.strip():
        raise ValueError("No text could be extracted by OCR engine")

    doc_date = detect_date(ocr_result.text)
    engine = ocr_result.engine or "unknown"
    page_count = ocr_result.metadata.get("page_count", len(ocr_result.pages))
    quality = "high" if page_count > 1 else "single_page"

    return {
        "text": ocr_result.text,
        "detected_date": doc_date,
        "source_type": ext.lstrip("."),
        "meta": {
            "engine": engine,
            "page_count": page_count,
            "pages": len(ocr_result.pages),
            "metadata": ocr_result.metadata,
        },
        "ocr_result": ocr_result,  # raw result for structured storage after doc_id is known
        "ocr_json": None,
        "ocr_markdown": None,
        "ocr_engine": engine,
        "page_count": page_count,
        "ocr_quality": quality,
    }
